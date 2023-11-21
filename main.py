# -*- coding: utf-8 -*-
"""
Created on Thu Mar 30 09:24:21 2023

@author: MEA8BU

"""

section = 0

# merge_Myc = True
merge_Myc = False

sectionAdd = True
#sectionAdd = False

selected= True
# selected = False

report2excel = True                                               # exports to excel & csv
#report2excel = False

prepare_slides = True
#prepare_slides = False

dataDesc = True
#dataDesc = False

# dataSelect = True # small
# dataSelect = False # big
# dataSelect = "all"
dataSelect = "new"

import sys
import os, time
import pandas as pd
import logging
import numpy as np
import matplotlib
import matplotlib.pyplot as plt
import seaborn as sns
import sklearn
import xgboost as xgb
import pickle
import warnings

warnings.filterwarnings("ignore", category=FutureWarning)
warnings.filterwarnings("ignore", category=DeprecationWarning)
from sklearn.model_selection import GridSearchCV, RandomizedSearchCV, cross_val_score, LeaveOneOut, KFold
from scipy.stats import randint, uniform, loguniform
from PIL import Image as PIL_Image
from datetime import datetime
from sklearn.preprocessing import StandardScaler
from sklearn.preprocessing import LabelEncoder
from sklearn.model_selection import train_test_split
from sklearn.model_selection import RepeatedStratifiedKFold
from pptx import Presentation
from pptx.util import Inches
# from win10toast import ToastNotifier

from sklearn.linear_model import LinearRegression, Ridge
from sklearn.tree import DecisionTreeRegressor
from sklearn.ensemble import RandomForestRegressor, GradientBoostingRegressor
from sklearn.neural_network import MLPRegressor
from sklearn.svm import SVR
from sklearn.neighbors import KNeighborsRegressor
from xgboost import XGBRegressor
from sklearn.metrics import mean_squared_error, mean_absolute_error, r2_score


if prepare_slides == True:
    import pptx
    ppt_template = "C:/Users/aysem/Desktop/Python/SalesForecast/Reports/Template.pptx"
    prs = Presentation(ppt_template)

RunID = datetime.now().strftime("%Y%m%d-%H%M")
logName = "log_" + str(RunID)
root = "C:/Users/aysem/Desktop/Python/SalesForecast/Reports"
pd.set_option("display.max_columns", None)

PIL_Image.MAX_IMAGE_PIXELS = None
sTime = time.time() #Return the current time in seconds .. Start time
# notifier = ToastNotifier()
data = 'data_new.xlsx'

#%%logging

def myLog(text, lev='info', end=None):
    print(text, end=end)
    encoded_text = text.encode('utf-8')
    logging.info(encoded_text)

log_file = root + RunID + '.log'

try:
    logging.basicConfig(filename=log_file,
                        format='%(asctime)s:\n%(message)s',
                        filemode='a',
                        level=logging.DEBUG,
                        )
    print("log oluşturuldu")
except:
    print("LOG OLUŞTURULAMADI")

#%% data and parameters

"""
different parameters for training the diffent datasets
"""

if selected == True:

    parameters = ['Tarih', 'Yıl', 'Müşteri Kodu', 'Ürün Kodu', 'Miktar', 'Birim Fiyat',
           'Toplam', 'Net Tutar']


#%% slide functions

"""
create powerpoint presentations for reporting
"""
def slide_addSection(sTitle, sText):
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.placeholders[0]
    title.text = sTitle
    subtitle = slide.placeholders[1]
    subtitle.text = sText
    return slide

def slide_addSlide(cTitle="AgiePGS v2", sTitle="", cText=""):
    if cText=="":
        title_slide_layout = prs.slide_layouts[16]
    else:
        title_slide_layout = prs.slide_layouts[8] # .. with Textbox (with bullets)
    slide = prs.slides.add_slide(title_slide_layout)
#    title = slide.shapes.title       # Slide Title .. also working
#    title.text = sTitle
    title = slide.placeholders[15]    # Chapter Title
    title.text = cTitle
    title = slide.placeholders[0]     # Slide Title
    title.text = sTitle
    try:
        subtitle = slide.placeholders[1]  # Content
        subtitle.text = cText
    except:
        # no text box
        pass
    return slide

def slide_addText(slide, sText, txt_left, txt_top, txt_width=1, txt_height=1):
    txBox = slide.shapes.add_textbox(Inches(txt_left), Inches(txt_top), Inches(txt_width), Inches(txt_height))
    tf = txBox.text_frame
    tf.text = sText

def slide_addHeader(slide, sText, txt_left, txt_top, txt_width=1, txt_height=1, text_color=(0, 123, 192), start_left=1, start_top=1):
    txBox = slide.shapes.add_textbox(Inches(txt_left), Inches(txt_top), Inches(txt_width), Inches(txt_height))
    tf = txBox.text_frame
    tf.text = sText

def slide_addPicture(slide, img_path, img_left, img_top, img_width=None, img_height=None):
    if  img_width != None:
        img_width = Inches(img_width)
    if  img_height != None:
        img_height = Inches(img_height)
    slide.shapes.add_picture(img_path, Inches(img_left), Inches(img_top), width=img_width, height=img_height)


#%% Comparision of algorithms
"""
visualizate to comparision of all algorithms
"""

def rComparison(sText, aNames, names, results, cText=""):
    fig, ax = plt.subplots()
    ax.bar(names, results)
    ax.set_ylim([0.5, 1.0])
    plt.xticks(rotation=45, ha="right")

    for i, v in enumerate(results):
        ax.text(i, v + 0.01, f"{v:.3f}", ha="center")

    img = root + aNames + RunID + sText + ".png"
    plt.savefig(img, format="png", bbox_inches="tight")
    # plt.show()
    plt.close()

    if prepare_slides == True:
        slide = slide_addSlide(sTitle=sText, cText=cText)
#        slide_addText(slide, sText, 0.20, 1.4)
        slide_addPicture(slide, img, 6.0, 2.1, img_width=5)

#%% Read data
"""
import data from dataset (excel)
"""

df = pd.DataFrame(pd.read_excel(data))
label_encoder = LabelEncoder()
df.drop(['Müşteri Kodu', 'Net Tutar'])

data['Tarih'] = pd.to_datetime(data['Tarih'])
df['Ay'] = df['Tarih'].dt.month
df['Year'] = df['Tarih'].dt.year
df['Ay'] = df['Tarih'].astype('float')
df['Year'] = df['Tarih'].astype('float')
data['Miktar'] = data['Miktar'].astype('float')

data['İrsaliye No'] = data['İrsaliye No'].astype('object')
data['Ürün Kodu'] = data['Ürün Kodu'].astype('object')

data['Birim Fiyat'] = data['Birim Fiyat'].astype('float')
data['Toplam'] = data['Toplam'].astype('float')
data['Net Tutar'] = data['Net Tutar'].astype('float')

df['Ürün Kodu'] = label_encoder.fit_transform(df['Ürün Kodu'])



['Tarih', 'Yıl', 'Ürün Kodu', 'Miktar', 'Birim Fiyat',
       'Toplam', ]



print(df.head(10))

# Bağımsız değişkenleri (X) ve hedef değişkeni (y) ayır
X = df.drop('Miktar', axis=1)
y = df['Miktar']

# Train ve test veri setlerini oluştur
# X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)


# Verileri eğitim ve test setlerine ayır
X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2)

# Verileri ölçeklendir
scaler = StandardScaler()
X_train = scaler.fit_transform(X_train)
X_test = scaler.transform(X_test)


#%% info slide

öMesaj = ("\tPython version:\t\t{}".format(sys.version)
        + "\n \tPandas version:\t\t{}".format(pd.__version__)
        + "\n \tNumPy version:\t\t{}".format(np.__version__)
        + "\n \tSeaborn version:\t\t{}".format(sns.__version__)
        + "\n \tMatplotlib version:\t{}".format(matplotlib.__version__)
        + "\n \tscikit-learn version:\t{}".format(sklearn.__version__)
        + "\n \txgboost version:\t\t{}".format(xgb.__version__)
    )

if prepare_slides == True:
    öMesaj += "\n\tpython-pptx version:\t{}". format(pptx.__version__)

öMesaj += "\n\n\tPython code:\t" + __file__ +\
          "\n\n\tlast modified:\t%s" % time.ctime(os.path.getmtime(__file__))
öMesaj = öMesaj.replace("|","\n\t\t|")

pd.set_option('display.float_format', lambda x: '{:.3f}'.format(x))   #Limiting floats output to 3 decimal points
myLog(öMesaj + "\n")

if dataDesc == True:
    slide = slide_addSlide(sTitle="Python Environment")
    slide_addText(slide, öMesaj, 0.28, 1.3)

#%% Data Overview

öMesaj = ""
öMesaj += "Data:\t" + logName
öMesaj += "\n(RunID:\t" + RunID + " )\n"
öMesaj += "\n\tData source:\t" + data
öMesaj += "\n\tState:\t\t" + str(time.ctime(os.path.getctime(root)))
öMesaj += "\n\n\tRows:\t\t{:,d}\n\tColumns:\t{}".format(X_train.shape[0], X_train.shape[1]).replace(",", "'")
öMesaj += "\n\n\t      XTrain dim: \t{:,d} x {:,d}".format(X_train.shape[0], X_train.shape[1]).replace(",", "'")
öMesaj += "\n\t      XTest dim: \t{:,d} x {:,d}".format(X_test.shape[0], X_test.shape[1]).replace(",", "'")
öMesaj += "\n\n\t      yTrain dim: \t {:,d} x 1".format(y_train.shape[0]).replace(",", "'")
öMesaj += "\n\t      yTest dim: \t {:,d} x 1".format(y_test.shape[0]).replace(",", "'")

if (prepare_slides == True):
    slide = slide_addSlide(sTitle= "Data Overview")
    slide_addText(slide, öMesaj, 0.20, 1.4)

myLog(str(section) +". " + öMesaj + "\n")

#%% ML FUNCTION


def evaluate_regression_models(csv_file_path, target_col_name, feature_col_names):
    """
    algorithm descriptions for information
    """
    global section

    descs = {
        "linearReg": "From the implementation point of view, this is just plain Ordinary Least Squares wrapped as a predictor object.",
        "ridgeReg": "Linear regression with L2 regularization.",
        "desicionTree": "A machine learning algorithm that creates a decision tree based on the input features to make predictions.",
        "randomForest": "An ensemble learning algorithm that creates multiple decision trees and combines their outputs to make predictions.",
        "gradientBoost": "An ensemble learning algorithm that creates multiple decision trees, where each subsequent tree is trained to correct the errors of the previous tree.",
        "mlpReg": "A neural network-based algorithm that can be used for regression problems.",
        "xgbReg": "An implementation of gradient boosting that uses a more efficient tree-building process.",
        "svm": "A machine learning algorithm that tries to find a hyperplane that best separates the classes in the data.",
        "knn": "A machine learning algorithm in the supervised learning category that uses the k-nearest neighbors approach to classify new data points based on their proximity to existing data points in the feature space."
    }

    models = [
        # Linear Regression
        (
            LinearRegression(),
            'Linear Regression',
            {}
        ),

        # Ridge Regression
        (
            Ridge(),
            'Ridge Regression',
            {'alpha': uniform(0.1, 10)}  # Daha küçük bir aralık
        ),

        # Decision Tree Regressor
        (
            DecisionTreeRegressor(random_state=42),
            'Decision Tree Regressor',
            {'max_depth': [3, 5, 10, None]}
        ),

        # Random Forest Regressor
        (
            RandomForestRegressor(),
            'Random Forest Regressor',
            {'n_estimators': randint(100, 1001), 'max_depth': [5, 10, 20, None]}
        ),

        # Gradient Boosting Regressor
        (
            GradientBoostingRegressor(),
            'Gradient Boosting Regressor',
            {'n_estimators': randint(100, 1001), 'learning_rate': uniform(0.01, 0.5 - 0.01)}
        ),

        # MLP Regressor
        (
            MLPRegressor(random_state=42, max_iter=1000),
            'MLP Regressor',
            {
                'hidden_layer_sizes': [(50, 50, 50), (50, 100, 50), (100,)],
                'activation': ['tanh', 'relu'],
                'solver': ['sgd', 'adam'],
                'alpha': uniform(0.0001, 0.001),  # Daha küçük bir aralık
                'learning_rate': ['constant', 'adaptive']
            }
        ),

        # K-Nearest Neighbors Regressor
        (
            KNeighborsRegressor(),
            'KNeighbors Regressor',
            {'n_neighbors': [3, 5, 7, 9, 11, 13, 15], 'weights': ['uniform', 'distance']}
        ),

        # XGBoost Regressor
        (
            XGBRegressor(),
            'XGBoost Regressor',
            {'n_estimators': randint(100, 501), 'max_depth': [5, 10, 20, None], 'learning_rate': uniform(0.01, 0.1)}  # Daha küçük bir aralık
        ),

        # Support Vector Regressor (SVR)
        (
            SVR(),
            'SVR',
            {'C': loguniform(0.1, 10), 'kernel': ['linear', 'rbf']}  # Daha küçük bir aralık ve logaritmik dağılım
        )
    ]

    # Result listeleri
    mse_results = []
    mae_results = []
    r2_results = []
    names = []

    # Modellerin işlenmesi
    for model, aNames, param_distributions in models:
        section += 1
        dText = descs.get(aNames, "")
        dKey = aNames

        if dText:
            if prepare_slides:
                slide = slide_addSection(str(section) + ". " + aNames, str(dText))
            myLog(str(section) + ". " + aNames, str(dText))

        if aNames == "linearReg":
            # model = LogisticRegression(random_state =42, max_iter = 1000)
            model.fit(X_train, y_train)
        else:

            #%%  Hyperparameters
            """
             predefined settings or configurations in a machine learning model that are determined before training begins. 
            """

            """
            the beginning                     ----------------------------------------------------------
            """
            if prepare_slides == True:
                slide = slide_addSlide(sTitle="Default Hyperparameters")
                slide_addText(slide, str(model.get_params()).replace("'","").replace(", ","\n"), 0.20, 1.4)
            myLog("\n\tDefault Hyperparameters:\n\t" + str(model.get_params()))

            """
            Randomized Search                 ----------------------------------------------------------
            """
            # myLog("before Random: " + str(best_params_random))
            random_search_iterations = 10
            random_search = RandomizedSearchCV(model, param_distributions, n_iter=random_search_iterations, cv=5, n_jobs=-1)
            random_search.fit(X_train, y_train)
            best_params_random = random_search.best_params_

            if prepare_slides == True:
                slide = slide_addSlide(sTitle="RandomizedSearchCV Hyperparameter Results")
                slide_addText(slide, str(best_params_random).replace("'","").replace(", ","\n"), 0.20, 1.4)
            myLog("\n\tafter RandomizedSearchCV:\n\t\t" + str(best_params_random))

            """
            Grid SearchCV                     ----------------------------------------------------------
            """
            # param_grid values for Grid Search
            param_grid = {}

            if 'n_estimators' in param_distributions:
                param_grid['n_estimators'] = [best_params_random['n_estimators']]

            if 'max_depth' in param_distributions:
                param_grid['max_depth'] = [best_params_random['max_depth']]

            if 'learning_rate' in param_distributions:
                param_grid['learning_rate'] = [best_params_random['learning_rate']]

            # Grid Search
            grid_search = GridSearchCV(model, param_grid, cv=5, n_jobs=-1)
            grid_search.fit(X_train, y_train)
            best_params_grid = grid_search.best_params_

            if prepare_slides == True:
                slide = slide_addSlide(sTitle="GridSearchCV Hyperparameter Results")
                slide_addText(slide, str(best_params_grid).replace("'","").replace(", ","\n"), 0.20, 1.4)
            myLog("\n\tafter GridSearchCV:\n\t\t" + str(best_params_grid))

            """
            Fit Model with best parameters    ----------------------------------------------------------
            """
            # En iyi modeli oluşturun ve eğitin
    #        best_model = GradientBoostingClassifier(**best_params_grid)
            model = model.set_params(**best_params_grid)
    #        best_model = model(**best_params_grid)     # .. why not ??
            model.fit(X_train, y_train)

            if prepare_slides == True:
                slide = slide_addSlide(sTitle=model.__class__.__name__ + " Final Hyperparameter")
                slide_addText(slide, str(model.get_params()).replace("'","").replace(", ","\n"), 0.20, 1.4)
            myLog("\n\tafter fit:\n\t" + str(best_params_grid))


    #%% Validation
        """
        cross-validation is a crucial technique for assessing a machine learning model's 
        generalization performance, reducing the risk of overfitting, and making informed decisions 
        about model selection and hyperparameter tuning.
        """

        #kfold validation-----------------------a type of cross validation-----------------
        kfold = KFold(n_splits=5, shuffle=True, random_state=42)
        scores = cross_val_score(model, X, y, cv=kfold)

        #Repeated kfold validation--------------a type of cross validation-----------------
        cv = RepeatedStratifiedKFold(n_splits=10, n_repeats=5, random_state=42)
        scoresr = cross_val_score(model, X, y, scoring='accuracy', cv=cv)

        cvMesaj = (" "
                 + "\n Cross-Validation Score:\t\t{:.2f}".format(scores.mean()) + "\tdata divided into 5 clusters and validated"
                 + "\n\n Repeated Cross-Validation Score:\t{:.2f}".format(scoresr.mean()) + "\tdata divided into 10 clusters and validated (repeated)"
                 + "\n Repeated Standard Deviation:\t{:.2f}".format(scoresr.std())
                 + "\n\n"
                 )

        #LeaveOneOut validation ----------------a type of cross validation----------------
        loo = LeaveOneOut()
        looScores = cross_val_score(model, X, y, cv=loo)
        cvMesaj += ("Cross-validation LOO score"
                 + "\nLeave One Out Results:\n" + str(looScores)
                 + "\nLeave One Out Mean Accuracy:\t" + "{:.3f}".format(looScores.mean()))
        myLog("\n" + cvMesaj)

        if (prepare_slides == True):

            slide = slide_addSlide(sTitle="Cross-validation score")
            slide_addText(slide, cvMesaj, 0.20, 1.4)


    #%% classification results
        """
        help quantify how well a model is performing its intended task
        """
        y_pred_grid = model.predict(X_test)
        y_pred_train = model.predict(X_train)
        mse = mean_squared_error(y_test, y_pred_grid)
        mae = mean_absolute_error(y_test, y_pred_grid)
        r2 = r2_score(y_test, y_pred_grid)

        # Sonuçları listelere ekle
        mse_results.append(mse)
        mae_results.append(mae)
        r2_results.append(r2)
        names.append(aNames)

#%% Save as pickle
        filename = str(aNames) + RunID + '.pkl'
        with open(filename, 'wb') as file:
            pickle.dump(model, file)

#%% importances and coefficients
        """
        importance helps us understand which features are most relevant in making accurate predictions
        """

        if hasattr(model, 'feature_importances_'):
            iMesaj = "Top 10 Feature Importances:\n\n"
            iMesaj += str(dKey) + " \n\n"
            importances = model.feature_importances_
            features = X.columns
            indices = np.argsort(importances)[::-1][:10]
            top_importances = importances[indices]
            top_features = features[indices]
            plt.figure(figsize=(10,6))
            plt.title("Top 10 Feature Importances\n\n")
            sns.barplot(x=top_importances, y=top_features, orient='h')
            img_imp = root + aNames + RunID +"imp"+".png"
            plt.savefig(img_imp, format="png", bbox_inches="tight")
            plt.close("all")
            for i in range(len(top_features)):
                iMesaj += f"{top_features[i]}: {top_importances[i]:.3f}" +"\n"

            if (prepare_slides == True):
                slide = slide_addSlide(sTitle="Top importances")
                slide_addText(slide, iMesaj, 0.20, 1.4)
                slide_addPicture(slide, img_imp, 6.0, 1.8, img_width=5)
            myLog(iMesaj + "\n")

        """
        coefficients help us understand how changes in predictor variables impact the target variable.        
        
        """
        if hasattr(model, 'coef_'):
            iMesaj = str(dKey) + "\n"
            coefs = model.coef_
            features = X.columns
            indices = np.argsort(coefs[0])[::-1][:10]
            top_coefs = coefs[0][indices]
            top_features = features[indices]
            plt.figure(figsize=(10,6))
            plt.title("Top 10 Coefficients\n\n")
            sns.barplot(x=top_coefs, y=top_features, orient='h')
            img_coe = root + aNames + RunID +"coef"+".png"
            plt.savefig(img_coe, format="png", bbox_inches="tight")
            plt.close("all")
            plt.show()
            for i in range(len(top_features)):
                iMesaj += f"{top_features[i]}:{top_coefs[i]:.3f}" +"\n"

            if (prepare_slides == True):
                slide = slide_addSlide(sTitle="Top 10 Coefficients")
                slide_addText(slide, iMesaj, 0.20, 1.4)
                slide_addPicture(slide, img_coe, 6.0, 1.8, img_width=5)
            myLog(iMesaj + "\n")

#%% Roc Curve

        """
         a graphical representation method used to evaluate the performance of binary classification models.
        """
        for i in range(len(names)):

            öMesaj = f"Model: {names[i]}\n\n"
            öMesaj+= "MSE:  \t"+ "{:,.3f}".format(mse_results[i]) + "\n"
            öMesaj+= "MAE: \t"+ "{:,.3f}".format(mae_results[i])+ "\n"
            öMesaj+= "R2:    \t"+ "{:,.3f}".format(r2_results[i])+ "\n"

            myLog("\n" + öMesaj)


        if (prepare_slides == True):
            slide = slide_addSlide(sTitle="Train & Test model")
            slide_addText(slide, öMesaj, 0.20, 1.4)
            slide_addPicture(slide, 6.0, 1.8, img_width=5)

        plt.close("all")
        myLog(str(section) +". " +öMesaj + "\n")

    #%% HeatMap

        """
        Heatmaps for classification are valuable tools for understanding the strengths and  weaknesses of a classification 
        model,identifying areas for improvement,and making informed decisions about model adjustments.
        

        messageTest,img_rate_test, imgTest = generate_confusion_matrix(y_test, y_pred_grid, 'Test', aNames, RunID, root, section)

        if prepare_slides:
            slide = slide_addSlide(sTitle=str(aNames) + " results " + "(TEST)")
            slide_addText(slide, " " + messageTest, 0.20, 1.4)
            slide_addPicture(slide, imgTest, 7.4, 0.3, img_width=3.9)
            slide_addPicture(slide, img_rate_test,  7.4, 3.3, img_width=3.9)  # Use imgTrain here, not messageTrain

        messageTrain,img_rate_train ,imgTrain = generate_confusion_matrix(y_train, y_pred_train, 'Train', aNames, RunID, root, section)

        messageTrain  # Combine the messages
        if prepare_slides:
            slide = slide_addSlide(sTitle=str(aNames) + " results " + "(TRAIN)")
            slide_addText(slide, " " + messageTrain, 0.20, 1.4)
            slide_addPicture(slide, imgTrain, 7.4, 0.3, img_width=3.9)
            slide_addPicture(slide, img_rate_train,  7.4, 3.3, img_width=3.9)  # Use imgTrain here, not messageTrain
"""
#%% Comparison of algorithms
    """
    compare all models to find the best model
    """
    #FIXME
    slide = slide_addSection("Comparison of algorithms", "We need to use the algorithm that has the best classification results. So visualize and compare them")
    rComparison("mse_results" , aNames, names, mse_results, "The accuracy score is a common metric used to evaluate the performance of a model in classification tasks. It measures the percentage of correct predictions made by the model across all classes.")
    rComparison("mae_results", aNames, names, mae_results, "The precision score is another crucial metric used in binary classification tasks. It measures the accuracy of positive predictions made by the model.")
    rComparison("r2_results"   , aNames, names, r2_results, "The recall score, also known as sensitivity or true positive rate (TPR), is another important metric used in binary classification tasks. It measures the ability of the model to correctly identify positive instances out of all the actual positive instances.")

#%% Other operations

def main():
    evaluate_regression_models(data, "Miktar", parameters)
    if prepare_slides == True:
        prs.save(root + RunID + " " + logName + "_slide.pptx")
    global sTime
    sTime = (time.time() - sTime) / 60
    myLog("\nfinish  ..  ;)\n\t.. run time: {:.2f} min\n\n".format(sTime))
    # notifier.show_toast(" AgiePGS Check", "finished .. run time: {:.2f}\n\n".format(sTime))

    logging.shutdown()

if __name__ == "__main__":
    main()

